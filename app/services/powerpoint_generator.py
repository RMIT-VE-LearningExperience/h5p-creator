from __future__ import annotations

import io
import hashlib
import re
import tempfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from urllib.parse import urljoin

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.parts.image import Image
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

from app.core.config import settings
from app.services import canvas_files, canvas_lms


class CourseDeckError(RuntimeError):
    pass


class SlideSpec(BaseModel):
    title: str = Field(default="", max_length=110)
    takeaway: str = Field(default="", max_length=220)
    bullets: list[str] = Field(default_factory=list, max_length=5)
    visual_cue: str = Field(default="", max_length=160)
    source_image_index: int | None = None


class DeckSpec(BaseModel):
    title: str = Field(default="Course Summary", max_length=120)
    subtitle: str = Field(default="", max_length=180)
    slides: list[SlideSpec] = Field(default_factory=list, min_length=1, max_length=40)


@dataclass
class CourseImage:
    index: int
    url: str
    alt: str
    page_title: str


@dataclass
class CourseSource:
    course_id: int
    course_name: str
    course_code: str
    modules: list[dict[str, Any]]
    pages: list[dict[str, str]]
    files: list[dict[str, str]]
    images: list[CourseImage]


async def build_course_source(
    course_id: int,
    include_images: bool = True,
    selected_module_ids: list[int] | None = None,
) -> CourseSource:
    course = await canvas_lms.read_course(course_id)
    course_info = course.get("course") or {}
    selected_ids = {int(module_id) for module_id in selected_module_ids or [] if module_id}
    all_modules = course.get("modules") or []
    modules = [
        module for module in all_modules
        if not selected_ids or int(module.get("id") or 0) in selected_ids
    ]
    module_page_urls, module_file_ids = _module_item_targets(modules)
    pages = []
    images: list[CourseImage] = []

    page_summaries = course.get("pages") or []
    if module_page_urls:
        page_summaries = [page for page in page_summaries if page.get("url") in module_page_urls]

    for page_summary in page_summaries[:25]:
        page_url = page_summary.get("url") or ""
        if not page_url:
            continue
        try:
            page = await canvas_lms.read_page(course_id, page_url)
        except canvas_lms.CanvasAPIError:
            continue
        text = _html_to_text(page.get("body", ""))
        if text:
            pages.append({
                "title": page.get("title") or page_summary.get("title") or page_url,
                "url": page_url,
                "text": _truncate(text, 6500),
            })
        if include_images:
            for img in _extract_images(page.get("body", ""), settings.canvas_base_url or ""):
                images.append(CourseImage(
                    index=len(images) + 1,
                    url=img["url"],
                    alt=img["alt"],
                    page_title=page.get("title") or page_url,
                ))

    files = []
    try:
        readable_files = await canvas_lms.list_course_files(course_id, limit=60)
    except canvas_lms.CanvasAPIError:
        readable_files = []
    if module_file_ids:
        readable_files = [item for item in readable_files if int(item.get("id") or 0) in module_file_ids]
    for item in readable_files[:8]:
        file_id = item.get("id")
        if not file_id:
            continue
        try:
            parsed = await canvas_files.parse_canvas_file(int(file_id))
        except Exception:
            continue
        if parsed.text.strip():
            files.append({
                "filename": parsed.filename,
                "title": parsed.title,
                "text": _truncate(parsed.text, 7000),
            })

    return CourseSource(
        course_id=course_id,
        course_name=course_info.get("name") or f"Course {course_id}",
        course_code=course_info.get("course_code") or "",
        modules=modules,
        pages=pages,
        files=files,
        images=images,
    )


def generate_deck_spec(
    source: CourseSource,
    slide_count: int,
    audience: str = "",
    include_images: bool = True,
    model_override: str | None = None,
) -> tuple[DeckSpec, str, int, int]:
    try:
        from openai import OpenAI
    except ImportError as exc:
        raise CourseDeckError("OpenAI SDK is not installed.") from exc

    if not settings.val_api_key:
        raise CourseDeckError("VAL_API_KEY is not configured")

    model_name = model_override or settings.val_model
    client = OpenAI(api_key=settings.val_api_key, base_url=settings.val_base_url)
    user_content = _build_prompt(source, slide_count, audience, include_images)

    try:
        response = client.chat.completions.create(
            model=model_name,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": _SYSTEM_PROMPT},
                {"role": "user", "content": user_content},
            ],
        )
    except Exception as exc:
        msg = str(exc)
        if any(x in msg for x in ("403", "Forbidden", "forbidden", "blocked", "unavailable")):
            raise CourseDeckError("VAL_NETWORK_ERROR") from exc
        raise

    raw = _strip_markdown_fences((response.choices[0].message.content or "").strip())
    if not raw:
        raise CourseDeckError("VAL returned no slide plan")
    spec = DeckSpec.model_validate_json(raw)
    spec.slides = spec.slides[: max(1, min(slide_count, 40))]
    if not spec.slides:
        raise CourseDeckError("VAL returned an empty slide plan")

    usage = getattr(response, "usage", None)
    return (
        spec,
        model_name,
        getattr(usage, "prompt_tokens", 0) or 0,
        getattr(usage, "completion_tokens", 0) or 0,
    )


async def render_deck(
    spec: DeckSpec,
    source: CourseSource,
    include_images: bool = True,
) -> bytes:
    prs = _new_presentation()

    _add_title_slide(prs, spec, source)
    image_cache = await _download_images(source.images) if include_images else {}
    image_paths = list(image_cache.values())
    used_image_paths: set[Path] = set()
    for idx, slide_spec in enumerate(spec.slides, start=1):
        image_path = _select_image_path(slide_spec, image_cache, image_paths, used_image_paths)
        layout_name = _choose_content_layout(idx, slide_spec, image_path)
        _add_content_slide(prs, slide_spec, idx, len(spec.slides), image_path, layout_name)
    _add_source_slide(prs, source)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


_SYSTEM_PROMPT = """\
You condense Canvas course content into a coherent PowerPoint slide plan for vocational education.

Return only JSON with this shape:
{
  "title": "Specific course deck title",
  "subtitle": "Short audience/context line",
  "slides": [
    {
      "title": "Slide title",
      "takeaway": "One-sentence main point",
      "bullets": ["3 to 5 concise bullets, not copied verbatim"],
      "visual_cue": "Brief image or diagram suggestion",
      "source_image_index": 1
    }
  ]
}

Rules:
- Condense and rephrase. Do not copy long passages verbatim.
- Preserve the logical course sequence and key module concepts.
- Prefer practical, learner-facing wording.
- Avoid filler slides such as "Introduction" unless the course content requires it.
- Use source_image_index only when a listed image is clearly relevant. Otherwise use null.
- Prefer using images on practical process, equipment, safety and standards slides when relevant images are available.
- Keep bullets short enough for a PowerPoint slide.
- If the audience/context line names specific required topics, give each one its own dedicated slide with a title that clearly names that topic. Do not fold a required topic into a generic summary slide alongside other required topics.
"""


def _build_prompt(source: CourseSource, slide_count: int, audience: str, include_images: bool) -> str:
    lines = [
        f"Course ID: {source.course_id}",
        f"Course name: {source.course_name}",
        f"Course code: {source.course_code}",
        f"Requested slide count: {slide_count}",
        f"Audience/context: {audience or 'RMIT vocational education learners'}",
        "",
        "## Modules",
    ]
    for module in source.modules[:30]:
        status = "published" if module.get("published") else "unpublished"
        lines.append(f"- {module.get('name', '')} ({module.get('items_count', 0)} items, {status})")

    if include_images and source.images:
        lines.extend(["", "## Available images"])
        for image in source.images:
            label = image.alt or "no alt text"
            lines.append(f"- {image.index}: {label} (from page: {image.page_title})")

    if source.pages:
        lines.extend(["", "## Canvas page text"])
        for page in source.pages:
            lines.extend([f"### {page['title']}", page["text"], ""])

    if source.files:
        lines.extend(["", "## Course file text"])
        for file in source.files:
            lines.extend([f"### {file['filename']}", file["text"], ""])

    return "\n".join(lines)


def _module_item_targets(modules: list[dict[str, Any]]) -> tuple[set[str], set[int]]:
    page_urls: set[str] = set()
    file_ids: set[int] = set()
    for module in modules:
        for item in module.get("items") or []:
            item_type = str(item.get("type") or "").lower()
            if item_type == "page" and item.get("page_url"):
                page_urls.add(str(item["page_url"]))
            if item_type == "file" and item.get("content_id"):
                try:
                    file_ids.add(int(item["content_id"]))
                except (TypeError, ValueError):
                    continue
    return page_urls, file_ids


def _add_title_slide(prs: Presentation, spec: DeckSpec, source: CourseSource) -> None:
    slide = prs.slides.add_slide(_layout_by_name(prs, "Title_option_1_text_only"))
    title = _placeholder(slide, 100) or slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(9.1), Inches(1.3))
    subtitle = _placeholder(slide, 101) or slide.shapes.add_textbox(Inches(0.3), Inches(1.85), Inches(9.1), Inches(1.0))
    _set_text(title, spec.title or source.course_name, Pt(34), bold=True)
    _set_text(subtitle, spec.subtitle or source.course_code or f"Canvas course {source.course_id}", Pt(18))


_DARK_BACKGROUND_LAYOUTS = {"Content with image_option 1", "Content with image_option 2", "Section divider_option 1"}
_IMAGE_SIDEBAR_LAYOUTS = {"Content with image_option 1", "Content with image_option 2"}
_TITLE_CHAR_WIDTH_FACTOR = 136.5  # empirical chars*pt per inch of column width, bold serif title font

# Layouts whose body placeholder defaults to nearly full slide width; when an image
# is added it must be narrowed so the generic image box (~6.05in-9.45in) doesn't
# overlap the text column.
_IMAGE_NARROWED_BODY: dict[str, dict[str, int]] = {
    "Content_option 1": {"left": Inches(0.42), "top": Inches(1.24), "width": Inches(5.35), "height": Inches(3.58)},
    "Content_option 2": {"left": Inches(3.02), "width": Inches(2.85)},
    "Content_option 3": {"left": Inches(3.02), "width": Inches(2.85)},
    "Section divider_option 1": {"left": Inches(3.02), "width": Inches(2.85)},
}


def _fit_sidebar_title_font(title: str, box_width_emu: int) -> tuple[int, int]:
    box_width_in = box_width_emu / 914400
    for font_pt in (24, 20, 18, 16):
        chars_per_line = max(4, int((_TITLE_CHAR_WIDTH_FACTOR * box_width_in) / font_pt))
        lines = max(1, -(-len(title) // chars_per_line)) if title else 1
        if lines <= 2 or font_pt == 16:
            return font_pt, lines
    return 16, 1


def _resize_sidebar_body(title_box: Any, body: Any, title_font_pt: int, title_lines: int) -> None:
    line_height_in = title_font_pt * 1.22 / 72
    needed_title_h_in = title_lines * line_height_in + 0.12
    title_top_in = title_box.top / 914400
    default_body_left = body.left
    default_body_width = body.width
    default_body_top_in = body.top / 914400
    default_body_bottom_in = (body.top + body.height) / 914400
    new_body_top_in = max(default_body_top_in, title_top_in + needed_title_h_in)
    if new_body_top_in > default_body_top_in:
        # Setting .top/.height materializes a fresh <a:xfrm> on placeholders that were
        # relying on inherited layout geometry; left/width must be re-asserted explicitly
        # or python-pptx defaults them to 0, collapsing the text column.
        body.left = default_body_left
        body.width = default_body_width
        body.top = Inches(new_body_top_in)
        body.height = Inches(max(0.5, default_body_bottom_in - new_body_top_in))


def _add_content_slide(
    prs: Presentation,
    spec: SlideSpec,
    index: int,
    total: int,
    image_path: Path | None,
    layout_name: str = "Content_option 1",
) -> None:
    slide = prs.slides.add_slide(_layout_by_name(prs, layout_name))

    has_image = image_path is not None
    title_box = _placeholder_by_type(slide, PP_PLACEHOLDER.TITLE) or slide.shapes.add_textbox(Inches(0.27), Inches(0.32), Inches(9.45), Inches(0.8))
    body = _placeholder_by_type(slide, PP_PLACEHOLDER.BODY) or slide.shapes.add_textbox(Inches(0.35), Inches(1.18), Inches(9.1), Inches(3.7))

    if has_image and layout_name in _IMAGE_NARROWED_BODY:
        # Placeholders inheriting geometry from the layout have no explicit <a:xfrm>;
        # writing any single dimension materializes one and defaults the untouched
        # dimensions to 0. Capture the inherited box first and always write all four
        # dimensions together so nothing collapses.
        base_left, base_top, base_width, base_height = body.left, body.top, body.width, body.height
        overrides = _IMAGE_NARROWED_BODY[layout_name]
        body.left = overrides.get("left", base_left)
        body.top = overrides.get("top", base_top)
        body.width = overrides.get("width", base_width)
        body.height = overrides.get("height", base_height)

    title_font_pt = 24
    if layout_name in _IMAGE_SIDEBAR_LAYOUTS:
        title_font_pt, title_lines = _fit_sidebar_title_font(spec.title, title_box.width)
        _resize_sidebar_body(title_box, body, title_font_pt, title_lines)
    _set_text(title_box, spec.title, Pt(title_font_pt), bold=True)

    dark_bg = layout_name in _DARK_BACKGROUND_LAYOUTS
    takeaway_color = RGBColor(255, 255, 255) if dark_bg else RGBColor(60, 60, 60)
    bullet_color = RGBColor(255, 255, 255) if dark_bg else RGBColor(20, 20, 20)
    cue_color = RGBColor(225, 225, 235) if dark_bg else RGBColor(90, 90, 90)

    frame = body.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = spec.takeaway
    p.font.size = Pt(13 if has_image else 15)
    p.font.bold = True
    p.font.color.rgb = takeaway_color
    p.space_after = Pt(10)

    bullets = [b.strip() for b in spec.bullets if b.strip()][:5]
    if not bullets and spec.visual_cue:
        bullets = [spec.visual_cue]
    for i, bullet in enumerate(bullets):
        p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(_body_font_size(layout_name, has_image, len(bullets)))
        p.font.color.rgb = bullet_color
        p.space_after = Pt(5 if has_image else 8)

    if image_path:
        left, top, width, height = _image_box(slide, layout_name)
        _add_image_fit(slide, image_path, left, top, width, height)
    elif spec.visual_cue:
        cue = slide.shapes.add_textbox(Inches(6.2), Inches(4.22), Inches(3.2), Inches(0.55))
        cue_frame = cue.text_frame
        cue_frame.clear()
        p = cue_frame.paragraphs[0]
        p.text = spec.visual_cue
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = cue_color


def _add_source_slide(prs: Presentation, source: CourseSource) -> None:
    slide = prs.slides.add_slide(_layout_by_name(prs, "End slide"))
    title = _placeholder(slide, 100) or slide.shapes.add_textbox(Inches(0.3), Inches(0.25), Inches(9.1), Inches(1.3))
    _set_text(title, "Thank you", Pt(34), bold=True)
    body = slide.shapes.add_textbox(Inches(0.35), Inches(1.75), Inches(8.6), Inches(1.9))
    frame = body.text_frame
    frame.clear()
    items = [
        f"{len(source.modules)} modules reviewed",
        f"{len(source.pages)} Canvas pages included",
        f"{len(source.files)} readable course files included",
        f"{len(source.images)} page images considered",
    ]
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)


def _new_presentation() -> Presentation:
    template_path = Path(settings.powerpoint_template_path)
    if template_path.exists():
        prs = Presentation(template_path)
        _clear_slides(prs)
        return prs
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs


def _clear_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _layout_by_name(prs: Presentation, name: str) -> Any:
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[6]


def _placeholder(slide: Any, idx: int) -> Any | None:
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def _placeholder_by_type(slide: Any, placeholder_type: PP_PLACEHOLDER) -> Any | None:
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type == placeholder_type:
            return placeholder
    return None


_CONTENT_LAYOUT_ROTATION = [
    "Content with image_option 1",
    "Content_option 2",
    "Content with image_option 2",
    "Content_option 3",
    "Content_option 1",
    "Section divider_option 1",
]


def _choose_content_layout(index: int, spec: SlideSpec, image_path: Path | None) -> str:
    return _CONTENT_LAYOUT_ROTATION[(index - 1) % len(_CONTENT_LAYOUT_ROTATION)]


def _body_font_size(layout_name: str, has_image: bool, bullet_count: int) -> int:
    if layout_name.startswith("Content with image"):
        return 9 if bullet_count >= 4 else 10
    if has_image:
        return 12
    return 15 if bullet_count >= 5 else 16


def _image_box(slide: Any, layout_name: str) -> tuple[int, int, int, int]:
    object_placeholder = _placeholder_by_type(slide, PP_PLACEHOLDER.OBJECT)
    if object_placeholder is not None:
        return (
            object_placeholder.left,
            object_placeholder.top,
            object_placeholder.width,
            object_placeholder.height,
        )
    if layout_name == "Content_option 1":
        return Inches(6.15), Inches(1.24), Inches(3.45), Inches(3.35)
    return Inches(6.05), Inches(1.2), Inches(3.4), Inches(3.35)


def _select_image_path(
    slide_spec: SlideSpec,
    image_cache: dict[int, Path],
    image_paths: list[Path],
    used_image_paths: set[Path],
) -> Path | None:
    if slide_spec.source_image_index:
        requested = image_cache.get(slide_spec.source_image_index)
        if requested and requested not in used_image_paths:
            used_image_paths.add(requested)
            return requested

    for candidate in image_paths:
        if candidate not in used_image_paths:
            used_image_paths.add(candidate)
            return candidate
    return None


def _set_text(shape: Any, text: str, size: Pt, bold: bool = False) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold


def _paint_background(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_accent_bar(slide: Any) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 30, 42)
    shape.line.fill.background()


def _add_slide_number(slide: Any, index: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.72), Inches(6.88), Inches(0.95), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"{index}/{total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(90, 90, 90)


def _add_image_fit(slide: Any, path: Path, left: int, top: int, width: int, height: int) -> None:
    try:
        image = Image.from_file(str(path))
        image_width, image_height = (int(value) for value in image.size)
        if image_width <= 0 or image_height <= 0:
            return
        scale = min(width / image_width, height / image_height)
        fitted_width = int(image_width * scale)
        fitted_height = int(image_height * scale)
        fitted_left = left + int((width - fitted_width) / 2)
        fitted_top = top + int((height - fitted_height) / 2)
        slide.shapes.add_picture(str(path), fitted_left, fitted_top, width=fitted_width, height=fitted_height)
    except Exception:
        return


async def _download_images(images: list[CourseImage]) -> dict[int, Path]:
    downloaded: dict[int, Path] = {}
    seen_hashes: set[str] = set()
    async with httpx.AsyncClient(headers=_canvas_headers(), timeout=20.0, follow_redirects=True) as client:
        for image in images:
            try:
                response = await client.get(image.url)
                if response.status_code >= 400:
                    continue
                digest = hashlib.sha1(response.content).hexdigest()
                if digest in seen_hashes:
                    continue
                seen_hashes.add(digest)
                suffix = _image_suffix(response.headers.get("content-type", ""), image.url)
                with tempfile.NamedTemporaryFile(prefix="course-img-", suffix=suffix, delete=False) as tmp:
                    tmp.write(response.content)
                    downloaded[image.index] = Path(tmp.name)
            except Exception:
                continue
    return downloaded


def _canvas_headers() -> dict[str, str]:
    if not settings.canvas_api_token:
        return {}
    return {"Authorization": f"Bearer {settings.canvas_api_token}"}


def _image_suffix(content_type: str, url: str) -> str:
    if "png" in content_type:
        return ".png"
    if "gif" in content_type:
        return ".gif"
    if "jpeg" in content_type or "jpg" in content_type:
        return ".jpg"
    match = re.search(r"\.(png|jpe?g|gif)(?:[?#]|$)", url, re.I)
    return f".{match.group(1).lower()}" if match else ".jpg"


class _HTMLExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.images: list[dict[str, str]] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attrs_dict = {key.lower(): value or "" for key, value in attrs}
        if tag in {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")
        if tag == "img" and attrs_dict.get("src"):
            self.images.append({
                "url": attrs_dict["src"],
                "alt": attrs_dict.get("alt", ""),
            })

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)


def _html_to_text(html: str) -> str:
    parser = _HTMLExtractor()
    parser.feed(html or "")
    text = " ".join(parser.parts)
    text = re.sub(r"\s*\n\s*", "\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _extract_images(html: str, base_url: str) -> list[dict[str, str]]:
    parser = _HTMLExtractor()
    parser.feed(html or "")
    images = []
    for image in parser.images:
        url = image["url"].replace("&amp;", "&").strip()
        if not url:
            continue
        images.append({
            "url": urljoin(base_url.rstrip("/") + "/", url),
            "alt": image.get("alt", "").strip(),
        })
    return images


def _truncate(value: str, max_chars: int) -> str:
    text = re.sub(r"\n{3,}", "\n\n", (value or "").strip())
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rstrip() + "\n[truncated]"


def _strip_markdown_fences(raw_json: str) -> str:
    if raw_json.startswith("```"):
        raw_json = raw_json.split("```")[1]
        if raw_json.startswith("json"):
            raw_json = raw_json[4:]
    return raw_json.strip()
