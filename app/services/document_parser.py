from __future__ import annotations
from dataclasses import dataclass, field
from pathlib import Path

from docx import Document
from docx.text.paragraph import Paragraph


@dataclass
class ParsedDocument:
    title: str
    sections: list[Section] = field(default_factory=list)
    raw_text: str = ""


@dataclass
class Section:
    heading: str
    level: int  # 1 = H1, 2 = H2, etc.
    paragraphs: list[str] = field(default_factory=list)
    # List items grouped: each sub-list is one bulleted/numbered list block
    lists: list[list[ListItem]] = field(default_factory=list)


@dataclass
class ListItem:
    text: str
    is_bold: bool = False  # bold = correct answer hint


def parse_docx(path: str | Path) -> ParsedDocument:
    doc = Document(str(path))
    title = ""
    sections: list[Section] = []
    current_section: Section | None = None
    current_list: list[ListItem] | None = None
    raw_lines: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            if current_list:
                if current_section:
                    current_section.lists.append(current_list)
                current_list = None
            continue

        raw_lines.append(text)
        style_name = _style_name(para)

        if style_name.startswith("heading"):
            if current_list and current_section:
                current_section.lists.append(current_list)
                current_list = None

            level = _heading_level(style_name)
            if level == 1 and not title:
                title = text
            current_section = Section(heading=text, level=level)
            sections.append(current_section)

        elif style_name in ("list paragraph", "list bullet", "list number") or _is_list(para):
            if current_section is None:
                current_section = Section(heading="", level=0)
                sections.append(current_section)
            if current_list is None:
                current_list = []
            bold = _has_bold(para)
            current_list.append(ListItem(text=text, is_bold=bold))

        else:
            if current_list and current_section:
                current_section.lists.append(current_list)
                current_list = None
            if current_section is None:
                current_section = Section(heading="", level=0)
                sections.append(current_section)
            current_section.paragraphs.append(text)

    if current_list and current_section:
        current_section.lists.append(current_list)

    if not title and sections:
        title = sections[0].heading or "Untitled Activity"

    return ParsedDocument(title=title, sections=sections, raw_text="\n".join(raw_lines))


def parse_pdf(path: str | Path) -> ParsedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValueError("pypdf is not installed. Run: pip install pypdf") from exc

    reader = PdfReader(str(path))
    raw_lines: list[str] = []

    for page in reader.pages:
        text = page.extract_text() or ""
        for line in text.split("\n"):
            line = line.strip()
            if line:
                raw_lines.append(line)

    if not raw_lines:
        return ParsedDocument(title="Empty PDF", sections=[], raw_text="")

    title = raw_lines[0]
    section = Section(heading=title, level=1, paragraphs=raw_lines[1:])
    return ParsedDocument(title=title, sections=[section], raw_text="\n".join(raw_lines))


def parse_pptx(path: str | Path) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("python-pptx is not installed. Run: pip install python-pptx") from exc

    prs = Presentation(str(path))
    sections: list[Section] = []
    raw_lines: list[str] = []
    title = ""

    for slide_idx, slide in enumerate(prs.slides):
        slide_title = ""
        body_texts: list[str] = []

        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                is_title = (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if is_title:
                        slide_title = text
                    else:
                        body_texts.append(text)
                    raw_lines.append(text)
            except Exception:
                continue

        if not slide_title:
            slide_title = f"Slide {slide_idx + 1}"
        if slide_idx == 0 and slide_title:
            title = slide_title

        sections.append(Section(heading=slide_title, level=2, paragraphs=body_texts))

    return ParsedDocument(
        title=title or "Untitled Presentation",
        sections=sections,
        raw_text="\n".join(raw_lines),
    )


def text_to_parsed_document(text: str, title: str = "Pasted Content") -> ParsedDocument:
    """Create a ParsedDocument from plain pasted text."""
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return ParsedDocument(title=title, sections=[], raw_text="")
    inferred_title = lines[0] if len(lines) > 1 else title
    section = Section(heading="", level=1, paragraphs=lines)
    return ParsedDocument(title=inferred_title, sections=[section], raw_text=text.strip())


def parse_any(path: str | Path) -> ParsedDocument:
    """Dispatch to the correct parser based on file extension."""
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return parse_docx(path)
    if ext == ".pdf":
        return parse_pdf(path)
    if ext == ".pptx":
        return parse_pptx(path)
    raise ValueError(f"Unsupported file type: {ext}. Use .docx, .pdf, or .pptx.")


def merge_documents(docs: list[ParsedDocument]) -> ParsedDocument:
    """Concatenate multiple parsed documents into one."""
    if not docs:
        return ParsedDocument(title="Empty", sections=[], raw_text="")
    if len(docs) == 1:
        return docs[0]
    merged_sections: list[Section] = []
    for doc in docs:
        merged_sections.extend(doc.sections)
    return ParsedDocument(
        title=docs[0].title,
        sections=merged_sections,
        raw_text="\n\n".join(doc.raw_text for doc in docs),
    )


def split_document(doc: ParsedDocument, n: int) -> list[ParsedDocument]:
    """Split a document into n roughly equal sub-documents by sections."""
    if n <= 1:
        return [doc]
    sections = doc.sections
    if not sections:
        return [doc] * n

    chunk_size = max(1, (len(sections) + n - 1) // n)
    result: list[ParsedDocument] = []
    for i in range(n):
        start = i * chunk_size
        chunk = sections[start : start + chunk_size] or sections
        raw = "\n".join(
            "\n".join(filter(None, [s.heading] + s.paragraphs)) for s in chunk
        )
        result.append(ParsedDocument(title=doc.title, sections=chunk, raw_text=raw))
    return result


# ── Internal helpers ───────────────────────────────────────────────────

def _heading_level(style_name: str) -> int:
    for i in range(1, 7):
        if str(i) in style_name:
            return i
    return 1


def _is_list(para: Paragraph) -> bool:
    return para._element.find(
        ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}numPr"
    ) is not None


def _has_bold(para: Paragraph) -> bool:
    return any(run.bold for run in para.runs if run.text.strip())


def _style_name(para: Paragraph) -> str:
    style = getattr(para, "style", None)
    name = getattr(style, "name", "") or ""
    return name.lower()
